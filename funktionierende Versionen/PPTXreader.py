import ast
import eurostat
import EStoPPTX as esi
from pptx import Presentation
#from pptx.enum.shapes import MSO_SHAPE
#from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
#from pptx.dml.color import RGBColor
from pptx.util import Inches, Cm
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

def description(my_chart:'object of the chart type of the pptx libary', show = False):
    '''Prints the chart type of the object, its Titel and whether it has a legend'''
    print('type: ', my_chart.chart_type)
    if my_chart.has_title:
        print('the chart has the titel: ', my_chart.title)
    if my_chart.has_legend:
        print('the chart has a Legend')
    if (my_chart.chart_type == XL_CHART_TYPE.LINE):
       print('Holding the following ',my_chart.series.count,' Lines')
       for L in my_chart.series:
        print (L.name)

if __name__ == "__main__":

    '''This is a test block to show the usage of the class
        goodpars = {"code": "STS_INPR_M",
                "flags": False,
                "indic_bt": ["PRD"],
                "nace_r2": ["C","D"],
                "s_adj": ["NSA"],
                "freq": ["M"],
                "unit": ["I21"],
                "geo": ["IT", "AT","DE"],
                "startPeriod": "2025-06",
                "endPeriod": "2025-07",
                "lang": "de",}
    '''
Update=False
prs = Presentation(pptx=r'D:\01 Nextcloud\Documents\Programmieren\Powerpoint Generator\Powerpoint\Output 2025-01.pptx')
print('slides:',len(prs.slides))

#print(len(prs.opc.coreprops))

'''This loop searches for slides holding Placeholders with charts.
If there is a chart, it checks the notes placeholder for a dataset code and a filter expression.
If they are found the Boolean 'Update' will be set to true and the Variable 'Chart'will be set
to the Chart which needs an update.
my_code holds the dataset-code and my_filter holds the filter expression to be used to retrieve the updated data.
For development purposes the function description prints some descriptions of the Chart.
'''
for sld in prs.slides:
    for pho in sld.placeholders:
        if(pho.has_chart):
            Chart = pho.chart

            '''
            hier musss das chart object gefunden und zugewiesens werden es solte ein
            Unteremelement des Grafikplatzhalters sein.
            '''
            if sld.has_notes_slide:
                    print('notes found on sld: ',sld.slide_id)
                    # print('text found for query',sld.notes_slide.notes_text_frame.text)
                    pars = dict()
                    for par in sld.notes_slide.notes_text_frame.paragraphs:
                        match par.text.split('=')[0].strip():
                            case 'my_lang':
                                my_lang = par.text.split('=')[1]
                                pars['lang']=my_lang.strip()
                                #print('LANG: ',my_lang)
                            case 'my_code':
                                my_code = par.text.split('=')[1]
                                pars['code']=my_code.strip()
                                #print('CODE: ',my_code)
                            case 'my_pars':
                                my_pars = eval((par.text.split('=')[1]))
                                pars.update(my_pars)
                                #print('FILTER: ',my_pars)
                            case 'my_time':
                                my_time = eval((par.text.split('=')[1]))
                                pars.update(my_time)
                                #print('TIME: ',my_time)
'''
Use the found parameters, to download the new data. Prepare the data to be inserted into the chart
 First select the metadata for the chosen dataset via get toc and put them into a dictonary
 use the fist list toc[0] containing the codes as key and the second list toc[1] containing
 the Meta data for my_code as description for the dataseries.


'''

print("gestartet")
print("aus dem File gezogen", pars)

# Create an object of the EStoPptxData class with the found parameters
T = esi.EStoPptxData(pars)
#T = esi.EStoPptxData(goodpars)
print(type(Chart))
description(Chart, show=True)

Chart.replace_data(T.ChartData)


prs.save(r'D:\01 Nextcloud\Documents\Programmieren\Powerpoint Generator\Powerpoint\Output 2025-02.pptx')



































