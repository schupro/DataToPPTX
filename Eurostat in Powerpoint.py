
import eurostat

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

'''
potential packages are:
pip install SDMXthon 
documented at: https://docs.sdmxthon.meaningfuldata.eu/index.html
pip install pandasdmx 
documented at: https://pandasdmx.readthedocs.io/en/v1.0/
pip install eurostat https://pypi.org/project/eurostat/
I give eurostat try 
'''
#toc_df = eurostat.get_toc_df() #read the full list of series offerd by eurostat into a dataframe (get_toc reads into a list)
#toc_df.to_excel("EurostatDatenangebot.xlsx")
#selection_df =eurostat.subset_toc_df(toc_df, 'hen') # read only series from toc_df where heading contains the Keyword given in parentheses

code = 'PRC_HICP_MANR'

parameters = eurostat.get_pars(code) #gives all parameters available in the Data given by code
print(parameters)
#values = eurostat.get_par_values(code, 'duration')
#values  
dictionary = eurostat.get_dic(code, 'coicop', full=False, frmt="dict", lang="de") #to get all values set full=True, instead of parameter name
#print(dictionary['N1-3'])
my_filter_pars = {'startPeriod': 2021,'endPeriod': 2022, 'geo': ['AT',], 'coicop': 'CP00'}
data = eurostat.get_data(code, filter_pars=my_filter_pars)
#dataInFrame = eurostat.get_data_df(code, filter_pars=my_filter_pars)
#data.to_excel("travel.xlsx")
 
# Öffnen der Präsentation mit Diagrammplatzhalter 
ppt = Presentation("Leere Vorlage.pptx")
new_slide = ppt.slides.add_slide(ppt.slide_layouts[21])
print (list(data[0][len(parameters):]))
print((str(data[1][:len(parameters)]), data[1][len(parameters):]))#this is the format needed for diagram data

# define chart data ---------------------
chart_data = CategoryChartData()
#chart_data.categories = ['East', 'West', 'Midwest']
#chart_data.add_series('Series 1', (19, 21, 16.7, 5))

chart_data.categories = list(data[0][len(parameters):]) 
chart_data.add_series(str(data[1][:len(parameters)]), data[1][len(parameters):])



# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
new_slide.placeholders[13].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
      
ppt.save("Output1.pptx")
# close the file               
#File_to_write_data.close()
  
print("Done")