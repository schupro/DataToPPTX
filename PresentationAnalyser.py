# Creating all layouts
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Cm

prs = Presentation("Leere Vorlage.pptx")

#Presentation Analyser, 
# Add a slide for every slide layout and 
# put the name and idx of the layout in the title

for layout in prs.slide_layouts:
    slide = prs.slides.add_slide(layout)
    if not(slide.shapes.title is None):
        slide.shapes.title.text = f"{layout.name} {prs.slide_layouts.index(layout)}"


# Analysing the shapes in each Slide:
prds=[]
for slide in prs.slides:
    if  not(slide.shapes.title is None):
        slds = [slide.slide_id, slide.slide_layout.name, prs.slide_layouts.index(layout)]
    else:
        slds = [slide.slide_id, "no Title"]
    shds=[] 
    for shape in slide.placeholders: 
        shds.append([shape.placeholder_format.idx, shape.left, shape.top, shape.name])
    prds.append([slds,shds])

#adding slide idx, Layout name, slide Layout number to notes, and shapes idx, and name
for slds in prds:
    sld = prs.slides.get(slds[0][0])
    sld.notes_slide.notes_text_frame.text = str(slds[0])
   
    for shds in slds[1]:        
        ind = sld.shapes.add_textbox(shds[1],shds[2], 5000, 5000)
        ind.fill.solid()
        ind.fill.fore_color.rgb = RGBColor(255,255,0)
        ind.text = str(shds[0])
        sld.notes_slide.notes_text_frame.add_paragraph().text = str(shds[0])+" "+str(shds[3])


"""
Ref for slide types:
0 ->  title and subtitle
1 ->  title and content
2 ->  section header
3 ->  two content
4 ->  Comparison
5 ->  Title only 
6 ->  Blank
7 ->  Content with caption
8 ->  Pic with caption
"""
prs.save("Analyse.pptx")
