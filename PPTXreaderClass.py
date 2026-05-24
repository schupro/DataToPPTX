import ast
import logging
import log_config
import eurostat
import EStoPPTX as esi
from pptx import Presentation
#from pptx.enum.shapes import MSO_SHAPE
#from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
#from pptx.dml.color import RGBColor
from pptx.util import Inches, Cm
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

logger = logging.getLogger(__name__)


class PPTXdataPresentation:
    def __init__(self, Filepath:str):
        self.prs = Presentation(pptx=Filepath)
        logger.info("PPTXdataPresentation Objekt erstellt mit Filepath: %s", Filepath)

    @property
    def slides(self):
        """PPTX slides iterator via property, so calling code can use prs.slides."""
        return self.prs.slides

    '''This loop searches for slides holding Placeholders with charts.
    If there is a chart, it checks the notes placeholder for a dataset code and a filter expression.
    If they are found the Boolean 'Update' will be set to true and the Variable 'Chart'will be set
    to the Chart which needs an update.
    my_code holds the dataset-code and my_filter holds the filter expression to be used to retrieve the updated data.
    For development purposes the function description prints some descriptions of the Chart.
    '''
    def update_charts(self):
        for sld in self.prs.slides:
            for pho in sld.placeholders:
                # print("ja ich bin hier")
                if pho.has_chart:
                    Chart = pho.chart
                    if sld.has_notes_slide:
                        logging.info('notes found on sld: %d', sld.slide_id)
                        logging.info('text found for query: %s', sld.notes_slide.notes_text_frame.text)
                        sld.pars = dict()
                        for par in sld.notes_slide.notes_text_frame.paragraphs:
                            match par.text.split('=')[0].strip():
                                case 'my_lang':
                                    my_lang =  par.text.split('=')[1]
                                    sld.pars['lang']=my_lang.strip()
                                    logging.info('LANG:  %s', my_lang)
                                case 'my_code':
                                    my_code = par.text.split('=')[1]
                                    sld.pars['code']=my_code.strip()
                                    logging.info('CODE:  %s', my_code)
                                case 'my_pars':
                                    my_pars = eval((par.text.split('=')[1]))
                                    sld.pars.update(my_pars)
                                    logging.info('FILTER:  %s', my_pars)
                                case 'my_time':
                                    my_time = eval((par.text.split('=')[1]))
                                    sld.pars.update(my_time)
                                    logging.info('TIME:  %s', my_time)

                        '''
                            Use the found parameters, to download the new data. Prepare the data to be inserted into the chart
                            First select the metadata for the chosen dataset via get toc and put them into a dictonary
                            use the fist list toc[0] containing the codes as key and the second list toc[1] containing
                            the Meta data for my_code as description for the dataseries.
                        '''


                        #logging.info("received parameters: %s", sld.pars)

                        # Create an object of the EStoPptxData class with the found parameters
                        T = esi.EStoPptxData(sld.pars)
                        #T = esi.EStoPptxData(goodpars)
                        # print(type(Chart))

                        #self.description(Chart, show=True)
                        Chart.replace_data(T.ChartData)


        def description(cls, my_chart:'object of the chart type of the pptx library', show = False) -> None:
            '''Prints the chart type of the object, its Titel and whether it has a legend'''
            print('type: ', my_chart.chart_type)
            if my_chart.has_title:
                print('the chart has the titel: ', my_chart.title)
            if my_chart.has_legend:
                print('the chart has a Legend')
            if (my_chart.chart_type == XL_CHART_TYPE.LINE):
                print('Holding the following ',my_chart.series.count,' Lines')
            for L in my_chart.series:
                print (L.name)


    def save(self, Filepath:str):
        self.prs.save(Filepath)


if __name__ == "__main__":

    log_config.set_log("logs/PPTXreaderClass.log")
    logger = logging.getLogger(__name__)
    logger.info("Projektlog von PPTXreaderClass main gestartet")


    '''This is a test block to show the usage of the class
        goodpars = {"code": "STS_INPR_M",
                "flags": False,
                "indic_bt": ["PRD"],
                "nace_r2": ["C","D"],
                "s_adj": ["NSA"],
                "freq": ["M"],
                "unit": ["I21"],
                "geo": ["IT", "AT","DE"],
                "startPeriod": "2025-06",
                "endPeriod": "2025-07",
                "lang": "de"}
    '''
    Update=False
    filename = r'D:\01 Nextcloud\Documents\Programmieren\Powerpoint Generator\Powerpoint\Output Test1.pptx'
    prs = PPTXdataPresentation(filename)
    logger.info ("reader gestartet%s", filename)
    logger.info('slides: %d', len(prs.slides))
    prs.update_charts()

    prs.save(filename)
