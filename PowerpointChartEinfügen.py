# import Presentation class
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
  
# Öffnen der Präsentation mit Diagrammplatzhalter 
ppt = Presentation("Leere Vorlage.pptx")
new_slide = ppt.slides.add_slide(ppt.slide_layouts[21])

# Beschreibung der Shapes in der Präsentation:
#for slide in ppt.slides: 
#    print("Folien ID: ", slide.slide_id, "Folien Name: ", slide.name, slide.shapes.title.text)
#    for shape in slide.placeholders: 
#        print('%d %s' % (shape.placeholder_format.idx, shape.placeholder_format.type))

#ppt.slides.get[1].placeholders[0].placeholder_format.idx
# Define chart data 
# Creating object of chart

# define chart data ---------------------
chart_data = CategoryChartData()
chart_data.categories = ['East', 'West', 'Midwest']
chart_data.add_series('Series 1', (19, 21, 16.7, 5))

# add chart to slide --------------------
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
new_slide.placeholders[13].insert_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, chart_data)
      
ppt.save("Output6.pptx")
# close the file               
#File_to_write_data.close()
  
print("Done")